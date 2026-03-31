"""
Text extraction for documents and presentations.
Returns (tag_text, full_text) tuples.
tag_text = first ~1000 words (sent to LLM for tagging)
full_text = entire extracted content (chunked for embedding)

Supported:
  .pdf              pdfplumber
  .docx             python-docx
  .pptx             python-pptx
  .xlsx             openpyxl
  .xls              xlrd
  .doc / .ppt       win32com (requires Microsoft Office on Windows)
"""

from pathlib import Path

WORDS_FOR_TAGGING = 1000
MAX_EXCEL_ROWS    = 500   # per sheet, to avoid embedding massive spreadsheets


def extract(path: Path) -> tuple[str, str]:
    """
    Returns (tag_text, full_text).
    Both strings are empty on failure — caller should skip the file.
    """
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _pdf(path)
        elif ext == ".docx":
            return _docx(path)
        elif ext == ".pptx":
            return _pptx(path)
        elif ext == ".xlsx":
            return _xlsx(path)
        elif ext == ".xls":
            return _xls(path)
        elif ext in (".doc", ".ppt"):
            return _com(path)   # requires Office on Windows
        else:
            return "", ""
    except Exception:
        return "", ""


# ── Helpers ───────────────────────────────────────────────────────────────────

def _trunc(text: str, words: int = WORDS_FOR_TAGGING) -> str:
    return " ".join(text.split()[:words])


def _pair(text: str) -> tuple[str, str]:
    return _trunc(text), text


# ── Extractors ────────────────────────────────────────────────────────────────

def _pdf(path: Path) -> tuple[str, str]:
    import pdfplumber
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
    return _pair("\n".join(pages))


def _docx(path: Path) -> tuple[str, str]:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _pair("\n".join(parts))


def _pptx(path: Path) -> tuple[str, str]:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return _pair("\n\n".join(slides))


def _xlsx(path: Path) -> tuple[str, str]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = [f"[Sheet: {ws.title}]"]
        count = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("\t".join(cells))
                count += 1
                if count >= MAX_EXCEL_ROWS:
                    rows.append(f"... (truncated at {MAX_EXCEL_ROWS} rows)")
                    break
        sheets.append("\n".join(rows))
    wb.close()
    return _pair("\n\n".join(sheets))


def _xls(path: Path) -> tuple[str, str]:
    import xlrd
    wb = xlrd.open_workbook(path)
    sheets = []
    for sheet in wb.sheets():
        rows = [f"[Sheet: {sheet.name}]"]
        for i in range(min(sheet.nrows, MAX_EXCEL_ROWS)):
            cells = [
                str(sheet.cell_value(i, j))
                for j in range(sheet.ncols)
                if str(sheet.cell_value(i, j)).strip()
            ]
            if cells:
                rows.append("\t".join(cells))
        sheets.append("\n".join(rows))
    return _pair("\n\n".join(sheets))


def _com(path: Path) -> tuple[str, str]:
    """
    Uses win32com to open .doc / .ppt via Microsoft Office (Windows only).
    Silently fails if Office is not installed.
    """
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
    except ImportError:
        return "", ""

    ext = path.suffix.lower()
    abs_path = str(path.resolve())
    text = ""

    try:
        if ext == ".doc":
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(abs_path, ReadOnly=True)
            text = doc.Content.Text
            doc.Close(False)
            word.Quit()
        elif ext == ".ppt":
            pp = win32com.client.Dispatch("PowerPoint.Application")
            prs = pp.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)
            parts = []
            for slide in prs.Slides:
                for shape in slide.Shapes:
                    try:
                        parts.append(shape.TextFrame.TextRange.Text)
                    except Exception:
                        pass
            text = "\n".join(parts)
            prs.Close()
            pp.Quit()
    except Exception:
        return "", ""
    finally:
        pythoncom.CoUninitialize()

    return _pair(text)
